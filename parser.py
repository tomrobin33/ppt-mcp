from pptx import Presentation
from typing import List, Dict, Any
from io import BytesIO
from docx import Document
import openpyxl
import tempfile
import os
from typing import Any, Dict
from PIL import Image
# PDF解析相关导入
import PyPDF2
import pdfplumber


def extract_text_from_shape(shape) -> List[str]:
    """
    从PPT形状中提取文本内容。
    支持以下内容提取：
    1. 普通文本框中的文本
    2. 表格中的文本
    3. 分组形状中的文本
    
    Args:
        shape: PPT中的形状对象
        
    Returns:
        包含所有提取文本的列表
    """
    texts = []
    if hasattr(shape, "text"):
        text = shape.text.strip()
        if text:
            texts.append(text)
    # 支持表格内容提取
    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
        for row in shape.table.rows:
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip()
                if cell_text:
                    texts.append(cell_text)
    # 支持分组 shape
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))
    return texts

def parse_pptx(file_bytes: bytes) -> Dict[str, Any]:
    """
    解析 PPTX 文件，返回结构化 JSON。
    
    功能说明：
    1. 支持内容：
       - 文本框中的文本
       - 表格中的文本
       - 分组形状中的文本
       
    2. 返回格式：
       {
           "slides": [
               {
                   "slide_index": 1,
                   "text": ["文本1", "文本2", ...]
               },
               ...
           ]
       }
    
    Args:
        file_bytes: PPTX文件的二进制内容
        
    Returns:
        包含所有幻灯片文本内容的字典
        
    Raises:
        ValueError: 当文件不是有效的PPTX格式时抛出
    """
    try:
        prs = Presentation(BytesIO(file_bytes))
    except Exception as e:
        raise ValueError(f"无法读取 pptx 文件: {e}")
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            try:
                texts.extend(extract_text_from_shape(shape))
            except Exception:
                continue
        slides.append({
            "slide_index": idx,
            "text": texts
        })
    return {"slides": slides}


def parse_docx(file_bytes: bytes) -> Dict[str, Any]:
    """
    解析 DOCX 文件，返回结构化 JSON。
    
    功能说明：
    1. 支持内容：
       - 文档中的所有段落文本
       - 表格内容（按行列结构保存）
       - 图片信息（文件名和大小）
       
    2. 返回格式：
       {
           "paragraphs": ["段落1", "段落2", ...],
           "tables": [
               [["单元格1", "单元格2"], ["单元格3", "单元格4"]],
               ...
           ],
           "images": [
               {"filename": "图片1.png", "size": 1024},
               ...
           ]
       }
    
    Args:
        file_bytes: DOCX文件的二进制内容
        
    Returns:
        包含文档内容的结构化字典
        
    注意：
    - 使用临时文件处理，会自动清理
    - 图片内容仅保存基本信息，不包含实际图片数据
    """
    result = {"paragraphs": [], "tables": [], "images": []}
    with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        doc = Document(tmp_path)
        # 段落
        for para in doc.paragraphs:
            result["paragraphs"].append(para.text)
        # 表格
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            result["tables"].append(table_data)
        # 图片
        rels = doc.part.rels
        for rel in rels:
            rel = rels[rel]
            if "image" in rel.target_ref:
                image_bytes = rel.target_part.blob
                result["images"].append({"filename": os.path.basename(rel.target_ref), "size": len(image_bytes)})
    finally:
        os.remove(tmp_path)
    return result


def parse_pdf(file_bytes: bytes) -> Dict[str, Any]:
    """
    解析 PDF 文件，返回结构化 JSON。
    
    功能说明：
    1. 支持内容：
       - 文档中的所有页面文本
       - 表格内容（使用pdfplumber提取）
       - 页面元数据（页数、页面大小等）
       - 图片信息（如果存在）
       
    2. 返回格式：
       {
           "pages": [
               {
                   "page_number": 1,
                   "text": "页面文本内容",
                   "tables": [
                       [["单元格1", "单元格2"], ["单元格3", "单元格4"]],
                       ...
                   ],
                   "images": [
                       {"bbox": [x1, y1, x2, y2], "type": "image"},
                       ...
                   ]
               },
               ...
           ],
           "metadata": {
               "total_pages": 10,
               "title": "文档标题",
               "author": "作者",
               "subject": "主题",
               "creator": "创建者"
           }
       }
    
    Args:
        file_bytes: PDF文件的二进制内容
        
    Returns:
        包含PDF内容的结构化字典
        
    Raises:
        ValueError: 当文件不是有效的PDF格式时抛出
    """
    result = {"pages": [], "metadata": {}}
    
    try:
        # 使用PyPDF2获取基本信息和文本
        pdf_reader = PyPDF2.PdfReader(BytesIO(file_bytes))
        
        # 获取元数据
        if pdf_reader.metadata:
            result["metadata"] = {
                "total_pages": len(pdf_reader.pages),
                "title": pdf_reader.metadata.get('/Title', ''),
                "author": pdf_reader.metadata.get('/Author', ''),
                "subject": pdf_reader.metadata.get('/Subject', ''),
                "creator": pdf_reader.metadata.get('/Creator', '')
            }
        else:
            result["metadata"] = {
                "total_pages": len(pdf_reader.pages),
                "title": "",
                "author": "",
                "subject": "",
                "creator": ""
            }
        
        # 使用pdfplumber进行更详细的解析
        with pdfplumber.open(BytesIO(file_bytes)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                page_data = {
                    "page_number": page_num,
                    "text": "",
                    "tables": [],
                    "images": []
                }
                
                # 提取文本
                page_text = page.extract_text()
                if page_text:
                    page_data["text"] = page_text.strip()
                
                # 提取表格
                tables = page.extract_tables()
                for table in tables:
                    if table:  # 确保表格不为空
                        page_data["tables"].append(table)
                
                # 提取图片信息
                images = page.images
                for img in images:
                    page_data["images"].append({
                        "bbox": img['bbox'],
                        "type": "image",
                        "width": img['width'],
                        "height": img['height']
                    })
                
                result["pages"].append(page_data)
                
    except Exception as e:
        raise ValueError(f"无法读取 PDF 文件: {e}")
    
    return result


def parse_xlsx(file_bytes: bytes) -> Dict[str, Any]:
    """
    解析 XLSX 文件，返回结构化 JSON。
    
    功能说明：
    1. 支持内容：
       - 所有工作表（sheets）的内容
       - 单元格的值和坐标信息
       - 单元格中的公式
       
    2. 返回格式：
       {
           "sheets": [
               {
                   "title": "Sheet1",
                   "cells": [
                       [
                           {"value": "A1的值", "coordinate": "A1"},
                           {"value": "B1的值", "coordinate": "B1"}
                       ],
                       ...
                   ],
                   "formulas": [
                       {"coordinate": "A1", "formula": "=SUM(B1:B10)"},
                       ...
                   ]
               },
               ...
           ]
       }
    
    Args:
        file_bytes: XLSX文件的二进制内容
        
    Returns:
        包含Excel文件内容的结构化字典
        
    注意：
    - 使用临时文件处理，会自动清理
    - data_only=False 设置可以获取公式内容
    """
    result = {"sheets": []}
    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
    try:
        wb = openpyxl.load_workbook(tmp_path, data_only=False)
        for sheet in wb.worksheets:
            sheet_data = {"title": sheet.title, "cells": [], "formulas": []}
            for row in sheet.iter_rows():
                row_data = []
                for cell in row:
                    cell_info = {"value": cell.value, "coordinate": cell.coordinate}
                    if cell.data_type == 'f':
                        cell_info["formula"] = cell.value
                        sheet_data["formulas"].append({"coordinate": cell.coordinate, "formula": cell.value})
                    row_data.append(cell_info)
                sheet_data["cells"].append(row_data)
            result["sheets"].append(sheet_data)
    finally:
        os.remove(tmp_path)
    return result